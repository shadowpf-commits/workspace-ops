#!/usr/bin/env python3
"""Gobs Pipeline PPT - PowerPoint version using python-pptx"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Colors
NAVY = RGBColor(0x1B, 0x36, 0x5D)
BLUE = RGBColor(0x44, 0x72, 0xC4)
LIGHT_BLUE = RGBColor(0x5B, 0x8B, 0xD4)
DARK_GRAY = RGBColor(0x2D, 0x37, 0x48)
MID_GRAY = RGBColor(0x71, 0x80, 0x96)
LIGHT_GRAY = RGBColor(0xE2, 0xE8, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_RED = RGBColor(0xC5, 0x30, 0x30)
GREEN = RGBColor(0x27, 0x67, 0x49)

def set_font(paragraph, size, bold=False, color=DARK_GRAY):
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color

def add_title_bar(slide, title_text, prs):
    """Add navy title bar at top"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, Inches(0.9)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()
    
    # Add title text
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    set_font(p, 24, bold=True, color=WHITE)

def add_footer(slide, prs):
    """Add footer bar"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(7.2),
        prs.slide_width, Inches(0.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()

def create_slide1_title(prs):
    """Slide 1: Title"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Navy header
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, Inches(2.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.8), Inches(12), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Gobs 功能 Pipeline"
    set_font(p, 44, bold=True, color=WHITE)
    
    # Subtitle
    p2 = tf.add_paragraph()
    p2.text = "AI 发行内容制作平台"
    set_font(p2, 22, color=LIGHT_BLUE)
    
    # Key message box
    shape2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.6), Inches(3),
        Inches(12), Inches(1)
    )
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = RGBColor(0xF7, 0xFA, 0xFC)
    shape2.line.color.rgb = LIGHT_GRAY
    
    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(3.2), Inches(11.6), Inches(0.7))
    tf2 = txBox2.text_frame
    p3 = tf2.paragraphs[0]
    p3.text = "核心目标：30天实现素材自动化 → 60天数据洞察闭环 → 年底完整独立发行"
    set_font(p3, 18, color=DARK_GRAY)
    
    # Date
    txBox3 = slide.shapes.add_textbox(Inches(0.6), Inches(4.5), Inches(12), Inches(0.4))
    tf3 = txBox3.text_frame
    p4 = tf3.paragraphs[0]
    p4.text = "2026年4月"
    set_font(p4, 12, color=MID_GRAY)
    
    add_footer(slide, prs)

def create_slide2_vision(prs):
    """Slide 2a: Vision"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "我们最终成为什么", prs)
    
    # Vision box
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.6), Inches(1.2),
        Inches(12), Inches(1)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()
    
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.6), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "一个基于 AI 的自动化发行平台：输入知识库，输出一键可执行的完整发行方案"
    set_font(p, 20, color=WHITE)
    
    # 3 pillars
    pillars = [
        ("01", "结构化知识库", "• 基于 Karpathy 工作流\n• 项目材料 → 冷启动\n• 知识积累"),
        ("02", "Agent 网络", "• 决策 / 执行 / 监控\n• 多类型 Agent 协同\n• 自动化工作流"),
        ("03", "一键输出", "• Community Plan\n• Budget/Event Plan\n• KOC 建设计划"),
    ]
    
    for i, (num, title, desc) in enumerate(pillars):
        x = 0.6 + i * 4.1
        
        # Card
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(2.5),
            Inches(3.8), Inches(3.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = LIGHT_GRAY
        
        # Number
        txBox = slide.shapes.add_textbox(Inches(x + 0.2), Inches(2.7), Inches(1), Inches(0.6))
        p = txBox.text_frame.paragraphs[0]
        p.text = num
        set_font(p, 28, bold=True, color=BLUE)
        
        # Title
        txBox2 = slide.shapes.add_textbox(Inches(x + 0.2), Inches(3.3), Inches(3.4), Inches(0.5))
        p2 = txBox2.text_frame.paragraphs[0]
        p2.text = title
        set_font(p2, 16, bold=True, color=NAVY)
        
        # Line
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x + 0.2), Inches(3.8),
            Inches(2), Inches(0.04)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = BLUE
        line.line.fill.background()
        
        # Description
        txBox3 = slide.shapes.add_textbox(Inches(x + 0.2), Inches(4), Inches(3.4), Inches(1.8))
        tf3 = txBox3.text_frame
        tf3.word_wrap = True
        p3 = tf3.paragraphs[0]
        p3.text = desc
        set_font(p3, 13, color=DARK_GRAY)
    
    add_footer(slide, prs)

def create_slide3_current(prs):
    """Slide 2b: Current Stage"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "我们目前在哪", prs)
    
    # Badge
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(8.5), Inches(0.25),
        Inches(1.2), Inches(0.4)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE
    shape.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(8.5), Inches(0.25), Inches(1.2), Inches(0.4))
    p = txBox.text_frame.paragraphs[0]
    p.text = "建设中"
    p.alignment = PP_ALIGN.CENTER
    set_font(p, 12, bold=True, color=WHITE)
    
    # 3 columns
    # Input
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(1.1),
        Inches(2.8), Inches(3.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = MID_GRAY
    
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(1.1),
        Inches(2.8), Inches(0.5)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = MID_GRAY
    header.line.fill.background()
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(2.8), Inches(0.5))
    p = txBox.text_frame.paragraphs[0]
    p.text = "输入 INPUT"
    p.alignment = PP_ALIGN.CENTER
    set_font(p, 12, bold=True, color=WHITE)
    
    txBox2 = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(2.4), Inches(2.3))
    tf = txBox2.text_frame
    tf.word_wrap = True
    items = ["• 游戏声量（爬虫）", "• 游戏舆情", "• Ingame数据", "• OM获客（暂无）"]
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        set_font(p, 12, color=DARK_GRAY)
    
    # Center - Platform
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(3.7), Inches(1.1),
        Inches(2.9), Inches(3.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()
    
    txBox = slide.shapes.add_textbox(Inches(3.7), Inches(1.8), Inches(2.9), Inches(1))
    p = txBox.text_frame.paragraphs[0]
    p.text = "Gobs"
    p.alignment = PP_ALIGN.CENTER
    set_font(p, 32, bold=True, color=WHITE)
    
    txBox2 = slide.shapes.add_textbox(Inches(3.7), Inches(2.6), Inches(2.9), Inches(0.5))
    p = txBox2.text_frame.paragraphs[0]
    p.text = "平台"
    p.alignment = PP_ALIGN.CENTER
    set_font(p, 14, color=LIGHT_BLUE)
    
    txBox3 = slide.shapes.add_textbox(Inches(3.7), Inches(3.2), Inches(2.9), Inches(0.5))
    p = txBox3.text_frame.paragraphs[0]
    p.text = "Campaign（规划中）"
    p.alignment = PP_ALIGN.CENTER
    set_font(p, 11, color=LIGHT_GRAY)
    
    # Output
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(6.9), Inches(1.1),
        Inches(2.8), Inches(3.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = GREEN
    shape.line.width = Pt(2)
    
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(6.9), Inches(1.1),
        Inches(2.8), Inches(0.5)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = GREEN
    header.line.fill.background()
    
    txBox = slide.shapes.add_textbox(Inches(6.9), Inches(1.1), Inches(2.8), Inches(0.5))
    p = txBox.text_frame.paragraphs[0]
    p.text = "输出 OUTPUT"
    p.alignment = PP_ALIGN.CENTER
    set_font(p, 12, bold=True, color=WHITE)
    
    txBox2 = slide.shapes.add_textbox(Inches(7.1), Inches(1.8), Inches(2.4), Inches(2.3))
    tf = txBox2.text_frame
    tf.word_wrap = True
    items = ["✓ 素材生成与分发", "✓ 水军运营", "", "→ Community Plan", "→ Budget Plan", "→ Event Plan", "→ KOC建设计划"]
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        set_font(p, 11, color=DARK_GRAY)
    
    # Note
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(4.5),
        Inches(12), Inches(0.6)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF7, 0xFA, 0xFC)
    shape.line.fill.background()
    
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(4.55), Inches(11.6), Inches(0.5))
    p = txBox.text_frame.paragraphs[0]
    p.text = "前提：进来需先绑定游戏 | 知识库基于 Andrej Karpathy 工作流"
    set_font(p, 11, color=MID_GRAY)
    
    add_footer(slide, prs)

def create_slide4_timeline(prs):
    """Slide 3: Timeline"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "路线图", prs)
    
    # Timeline phases
    phases = [
        ("现在", "4月底", "验证：AOV\n水军投入使用", ["• 素材生成", "• 水军运营", "• 视频粗剪"], MID_GRAY),
        ("4月底", "6月底", "验证：GNG泰国\n独立发行", ["• 数据洞察", "• 主动式Agent", "• 舆情监控"], BLUE),
        ("6月底", "年底", "验证：孵化项目\n全流程验证", ["• 完整闭环", "• 一键生成", "• 甩脱依赖"], NAVY),
    ]
    
    # Timeline line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.5), Inches(1.6),
        Inches(10.3), Inches(0.06)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = LIGHT_GRAY
    line.line.fill.background()
    
    for i, (label, date, verify, items, color) in enumerate(phases):
        x = 0.5 + i * 4.1
        
        # Dot
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x + 1.2), Inches(1.45),
            Inches(0.4), Inches(0.4)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = color
        dot.line.fill.background()
        
        # Label
        txBox = slide.shapes.add_textbox(Inches(x), Inches(1.1), Inches(2.8), Inches(0.3))
        p = txBox.text_frame.paragraphs[0]
        p.text = label
        p.alignment = PP_ALIGN.CENTER
        set_font(p, 11, bold=True, color=color)
        
        # Card
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(2),
            Inches(3.8), Inches(3.8)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = color
        shape.line.width = Pt(2)
        
        # Header
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(2),
            Inches(3.8), Inches(0.5)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = color
        header.line.fill.background()
        
        txBox = slide.shapes.add_textbox(Inches(x), Inches(2), Inches(3.8), Inches(0.5))
        p = txBox.text_frame.paragraphs[0]
        p.text = date
        p.alignment = PP_ALIGN.CENTER
        set_font(p, 14, bold=True, color=WHITE)
        
        # Verify
        txBox = slide.shapes.add_textbox(Inches(x + 0.15), Inches(2.6), Inches(3.5), Inches(0.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = verify.replace("\n", " ")
        set_font(p, 10, color=DARK_GRAY)
        
        # Items
        txBox2 = slide.shapes.add_textbox(Inches(x + 0.15), Inches(3.4), Inches(3.5), Inches(2.2))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        for j, item in enumerate(items):
            p = tf2.paragraphs[0] if j == 0 else tf2.add_paragraph()
            p.text = item
            set_font(p, 12, color=DARK_GRAY)
    
    add_footer(slide, prs)

def create_slide5_day30(prs):
    """Slide 4: Day 0-30"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Day 0-30：当前阶段进度（4月底目标）", prs)
    
    # Completed column
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(1.1),
        Inches(4.3), Inches(4.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = GREEN
    shape.line.width = Pt(2)
    
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(1.1),
        Inches(4.3), Inches(0.5)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = GREEN
    header.line.fill.background()
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(4.3), Inches(0.5))
    p = txBox.text_frame.paragraphs[0]
    p.text = "已完成 — 前20天"
    p.alignment = PP_ALIGN.CENTER
    set_font(p, 14, bold=True, color=WHITE)
    
    completed = ["视频内容分析", "AI 粗剪流程", "背景音乐匹配", "素材分发基础框架", "水军运营策略设计", "用户资产隔离"]
    txBox2 = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(3.9), Inches(3.6))
    tf = txBox2.text_frame
    tf.word_wrap = True
    for i, item in enumerate(completed):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "✓ " + item
        set_font(p, 13, color=DARK_GRAY)
    
    # In progress column
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(5.2), Inches(1.1),
        Inches(4.3), Inches(4.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = BLUE
    shape.line.width = Pt(2)
    
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(5.2), Inches(1.1),
        Inches(4.3), Inches(0.5)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = BLUE
    header.line.fill.background()
    
    txBox = slide.shapes.add_textbox(Inches(5.2), Inches(1.1), Inches(4.3), Inches(0.5))
    p = txBox.text_frame.paragraphs[0]
    p.text = "进行中 — 到4月底"
    p.alignment = PP_ALIGN.CENTER
    set_font(p, 14, bold=True, color=WHITE)
    
    remaining = ["视频生成标准化", "角色/场景一致性", "水军 Agent 上线", "系统稳定性测试", "FFmpeg 预处理优化"]
    txBox2 = slide.shapes.add_textbox(Inches(5.4), Inches(1.8), Inches(3.9), Inches(3.6))
    tf = txBox2.text_frame
    tf.word_wrap = True
    for i, item in enumerate(remaining):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "□ " + item
        set_font(p, 13, color=DARK_GRAY)
    
    add_footer(slide, prs)

def create_slide6_day60(prs):
    """Slide 5: Day 30-60"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "Day 30-60：6月底目标", prs)
    
    # Badge
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(8.3), Inches(0.25),
        Inches(1.2), Inches(0.4)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLUE
    shape.line.fill.background()
    txBox = slide.shapes.add_textbox(Inches(8.3), Inches(0.25), Inches(1.2), Inches(0.4))
    p = txBox.text_frame.paragraphs[0]
    p.text = "本期重点"
    p.alignment = PP_ALIGN.CENTER
    set_font(p, 11, bold=True, color=WHITE)
    
    # 3 parts
    parts = [
        ("数据洞察", ["• 爬虫接入", "• 声量分析", "• 舆情监控", "• ST API数据"]),
        ("主动式 Agent", ["• 发行方案生成", "• 决策链路", "• 执行与监控"]),
        ("输出能力", ["• Community Plan", "• Budget Plan", "• Event Plan", "• KOC建设"]),
    ]
    
    for i, (title, items) in enumerate(parts):
        x = 0.5 + i * 4.1
        
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(1.1),
            Inches(3.8), Inches(3)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = BLUE
        shape.line.width = Pt(2)
        
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(1.1),
            Inches(3.8), Inches(0.5)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = BLUE
        header.line.fill.background()
        
        txBox = slide.shapes.add_textbox(Inches(x), Inches(1.1), Inches(3.8), Inches(0.5))
        p = txBox.text_frame.paragraphs[0]
        p.text = title
        p.alignment = PP_ALIGN.CENTER
        set_font(p, 13, bold=True, color=WHITE)
        
        txBox2 = slide.shapes.add_textbox(Inches(x + 0.15), Inches(1.8), Inches(3.5), Inches(2.2))
        tf = txBox2.text_frame
        tf.word_wrap = True
        for j, item in enumerate(items):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = item
            set_font(p, 12, color=DARK_GRAY)
    
    # Note
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(4.3),
        Inches(12), Inches(0.6)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF7, 0xFA, 0xFC)
    shape.line.fill.background()
    
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(4.35), Inches(11.6), Inches(0.5))
    p = txBox.text_frame.paragraphs[0]
    p.text = "本期不重点：视频产出、水军运营（已上线，持续运行）"
    set_font(p, 12, color=MID_GRAY)
    
    add_footer(slide, prs)

def create_slide7_year_plan(prs):
    """Slide 6: Year plan"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "一年计划：90天详细路线（到6月底）", prs)
    
    plans = [
        ("01", "结构化知识库", "• Karpathy工作流\n• GNG材料冷启动", "进行中", BLUE),
        ("02", "GNG发行知识库", "• OM创意/预算\n• Community/Influencer\n• 方案质量≥80%", "规划中", NAVY),
        ("03", "更多Action Agent", "• Community Post Agent\n• 日报/周报Agent\n• Latte/社交/ST API", "规划中", NAVY),
    ]
    
    for i, (num, title, desc, status, color) in enumerate(plans):
        x = 0.5 + i * 4.1
        
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(1.1),
            Inches(3.8), Inches(3.2)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = color
        shape.line.width = Pt(2)
        
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(1.1),
            Inches(3.8), Inches(0.5)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = color
        header.line.fill.background()
        
        txBox = slide.shapes.add_textbox(Inches(x), Inches(1.1), Inches(3.8), Inches(0.5))
        p = txBox.text_frame.paragraphs[0]
        p.text = num + " " + title
        p.alignment = PP_ALIGN.CENTER
        set_font(p, 12, bold=True, color=WHITE)
        
        txBox2 = slide.shapes.add_textbox(Inches(x + 0.15), Inches(1.8), Inches(3.5), Inches(1.8))
        tf = txBox2.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc.replace("\n", " ")
        set_font(p, 12, color=DARK_GRAY)
        
        txBox3 = slide.shapes.add_textbox(Inches(x + 0.15), Inches(3.7), Inches(3.5), Inches(0.4))
        p = txBox3.text_frame.paragraphs[0]
        p.text = status
        set_font(p, 10, bold=True, color=color)
    
    # Note
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(4.5),
        Inches(12), Inches(0.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF7, 0xFA, 0xFC)
    shape.line.fill.background()
    
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(4.55), Inches(11.6), Inches(0.4))
    p = txBox.text_frame.paragraphs[0]
    p.text = "Agent List 扩展：覆盖规划/决策/执行/监控关键节点"
    set_font(p, 12, color=MID_GRAY)
    
    add_footer(slide, prs)

def create_slide8_year_end(prs):
    """Slide 7: Year end"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "年底目标", prs)
    
    # Two pillars
    pillars = [
        ("功能完成", ["• 完成 P1/P2 功能", "• 发行能力更健壮", "• Skills + MD 复用到更多产品"], NAVY),
        ("AI 优化思考", ["• 用户离研发更近：智能客服", "• 研发离用户更近", "  Agent 美化问题收集反馈"], MID_GRAY),
    ]
    
    for i, (title, items, color) in enumerate(pillars):
        x = 0.5 + i * 4.5
        
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(1.1),
            Inches(4.3), Inches(2.5)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = color
        shape.line.width = Pt(2)
        
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(1.1),
            Inches(4.3), Inches(0.5)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = color
        header.line.fill.background()
        
        txBox = slide.shapes.add_textbox(Inches(x), Inches(1.1), Inches(4.3), Inches(0.5))
        p = txBox.text_frame.paragraphs[0]
        p.text = title
        p.alignment = PP_ALIGN.CENTER
        set_font(p, 14, bold=True, color=WHITE)
        
        txBox2 = slide.shapes.add_textbox(Inches(x + 0.2), Inches(1.8), Inches(3.9), Inches(1.6))
        tf = txBox2.text_frame
        tf.word_wrap = True
        for j, item in enumerate(items):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = item
            set_font(p, 13, color=DARK_GRAY)
    
    # Incubator project
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(3.8),
        Inches(12), Inches(0.9)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()
    
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(3.9), Inches(11.6), Inches(0.7))
    p = txBox.text_frame.paragraphs[0]
    p.text = "孵化项目：一个产品全流程独立发行验证"
    set_font(p, 18, bold=True, color=WHITE)
    
    add_footer(slide, prs)

def create_slide9_actions(prs):
    """Slide 8: Actions"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_bar(slide, "立即行动项", prs)
    
    actions = [
        ("P0", "视频生成标准化 + 角色一致性", ACCENT_RED),
        ("P0", "水军 Agent 训练 + 上线", ACCENT_RED),
        ("P1", "爬虫接入 + 声量分析", BLUE),
        ("P1", "舆情监控分渠道配置", BLUE),
        ("P2", "系统稳定性测试（5-10用户）", MID_GRAY),
        ("P2", "用户空间权限隔离", MID_GRAY),
    ]
    
    for i, (priority, text, color) in enumerate(actions):
        y = 1.15 + i * 0.7
        
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(y),
            Inches(1), Inches(0.55)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(1), Inches(0.55))
        p = txBox.text_frame.paragraphs[0]
        p.text = priority
        p.alignment = PP_ALIGN.CENTER
        set_font(p, 12, bold=True, color=WHITE)
        
        txBox2 = slide.shapes.add_textbox(Inches(1.7), Inches(y), Inches(7.8), Inches(0.55))
        p = txBox2.text_frame.paragraphs[0]
        p.text = text
        set_font(p, 14, color=DARK_GRAY)
    
    # Goal
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(5.4),
        Inches(12), Inches(0.6)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()
    
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(5.45), Inches(11.6), Inches(0.5))
    p = txBox.text_frame.paragraphs[0]
    p.text = "4月底 AOV验证 → 6月底 GNG泰国独立发行 → 年底 完整独立发行"
    set_font(p, 13, bold=True, color=WHITE)
    
    add_footer(slide, prs)

def main():
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Create slides
    create_slide1_title(prs)
    create_slide2_vision(prs)
    create_slide3_current(prs)
    create_slide4_timeline(prs)
    create_slide5_day30(prs)
    create_slide6_day60(prs)
    create_slide7_year_plan(prs)
    create_slide8_year_end(prs)
    create_slide9_actions(prs)
    
    # Save
    output_path = "/Users/shadow/.openclaw/workspace/Gobs_Pipeline_PowerPoint.pptx"
    prs.save(output_path)
    print(f"✅ Done! {output_path}")

if __name__ == "__main__":
    main()
